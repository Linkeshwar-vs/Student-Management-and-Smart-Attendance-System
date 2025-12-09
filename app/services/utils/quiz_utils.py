# utils/quiz_utils.py
import os
from pathlib import Path
import fitz  # pip install pymupdf
import docx2txt  # pip install docx2txt
from pptx import Presentation  # pip install python-pptx
import textract  # optional fallback; pip install textract

def extract_text_from_file(filepath: str) -> str:
    ext = Path(filepath).suffix.lower()
    if ext == ".pdf":
        text = []
        doc = fitz.open(filepath)
        for page in doc:
            text.append(page.get_text("text"))
        return "\n\n".join(text)
    elif ext == ".docx":
        return docx2txt.process(filepath)
    elif ext in (".pptx", ".ppt"):
        prs = Presentation(filepath)
        slides_text = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, "text"):
                    slides_text.append(shp.text)
        return "\n\n".join(slides_text)
    else:
        # fallback: read as text
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception:
            # optional: try textract for odd types
            try:
                return textract.process(filepath).decode('utf-8')
            except Exception:
                return ""
